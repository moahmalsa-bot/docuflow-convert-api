from pathlib import Path

import fitz
import pdfplumber
from fastapi import HTTPException
from openpyxl import Workbook
from pdf2docx import Converter
from PIL import Image
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter

from app.utils.files import ConversionError, ghostscript_binary, libreoffice_binary, run_command, zip_directory


def pdf_has_selectable_text(pdf_path: Path, min_chars: int = 20) -> bool:
    try:
        with fitz.open(pdf_path) as document:
            for page in document:
                if len(page.get_text("text").strip()) >= min_chars:
                    return True
    except Exception:
        return False
    return False


def ocr_pdf(input_pdf: Path, output_pdf: Path) -> Path:
    run_command(
        [
            "ocrmypdf",
            "--skip-text",
            "--deskew",
            "--rotate-pages",
            "--optimize",
            "1",
            str(input_pdf),
            str(output_pdf),
        ],
        timeout=900,
    )
    return output_pdf


def pdf_to_word(input_pdf: Path, output_docx: Path, use_ocr: bool = False) -> Path:
    source_pdf = input_pdf
    if use_ocr or not pdf_has_selectable_text(input_pdf):
        source_pdf = input_pdf.with_name(f"{input_pdf.stem}.ocr.pdf")
        ocr_pdf(input_pdf, source_pdf)

    converter = Converter(str(source_pdf))
    try:
        converter.convert(str(output_docx), start=0, end=None)
    finally:
        converter.close()
    return output_docx


def pdf_to_powerpoint(input_pdf: Path, output_pptx: Path, dpi: int = 180) -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    blank_layout = presentation.slide_layouts[6]
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)

    with fitz.open(input_pdf) as document:
        for index, page in enumerate(document, start=1):
            image_path = output_pptx.with_name(f"page-{index:04d}.png")
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            pixmap.save(image_path)
            slide = presentation.slides.add_slide(blank_layout)
            slide.shapes.add_picture(
                str(image_path),
                0,
                0,
                width=presentation.slide_width,
                height=presentation.slide_height,
            )

    if len(presentation.slides) == 0:
        raise HTTPException(status_code=400, detail="PDF has no pages")
    presentation.save(output_pptx)
    return output_pptx


def pdf_to_excel(input_pdf: Path, output_xlsx: Path) -> Path:
    workbook = Workbook()
    default_sheet = workbook.active
    default_sheet.title = "Page 1"
    wrote_any = False

    with pdfplumber.open(input_pdf) as pdf:
        if not pdf.pages:
            raise HTTPException(status_code=400, detail="PDF has no pages")

        for page_index, page in enumerate(pdf.pages, start=1):
            sheet = default_sheet if page_index == 1 else workbook.create_sheet(title=f"Page {page_index}")
            tables = page.extract_tables()
            row_cursor = 1

            if tables:
                for table in tables:
                    for row in table:
                        for col_index, value in enumerate(row or [], start=1):
                            sheet.cell(row=row_cursor, column=col_index, value=value)
                        row_cursor += 1
                    row_cursor += 1
                    wrote_any = True
            else:
                text = page.extract_text() or ""
                for line in text.splitlines():
                    sheet.cell(row=row_cursor, column=1, value=line)
                    row_cursor += 1
                    wrote_any = True

            for column_cells in sheet.columns:
                values = [str(cell.value) for cell in column_cells if cell.value is not None]
                if values:
                    sheet.column_dimensions[column_cells[0].column_letter].width = min(max(map(len, values)) + 2, 80)

    if not wrote_any:
        default_sheet.cell(row=1, column=1, value="No tables or text found in PDF.")
    workbook.save(output_xlsx)
    return output_xlsx


def pdf_to_images_zip(input_pdf: Path, output_zip: Path, image_format: str, dpi: int = 200) -> Path:
    if dpi < 50 or dpi > 600:
        raise HTTPException(status_code=400, detail="dpi must be between 50 and 600")

    ext = "jpg" if image_format.lower() in {"jpg", "jpeg"} else "png"
    zoom = dpi / 72
    matrix = fitz.Matrix(zoom, zoom)
    files: list[Path] = []

    with fitz.open(input_pdf) as document:
        if document.page_count == 0:
            raise HTTPException(status_code=400, detail="PDF has no pages")
        for index, page in enumerate(document, start=1):
            pixmap = page.get_pixmap(matrix=matrix, alpha=False)
            image_path = output_zip.with_name(f"page-{index:04d}.{ext}")
            if ext == "jpg":
                temp_png = image_path.with_suffix(".tmp.png")
                pixmap.save(temp_png)
                with Image.open(temp_png) as image:
                    image.convert("RGB").save(image_path, "JPEG", quality=95, optimize=True)
                temp_png.unlink(missing_ok=True)
            else:
                pixmap.save(image_path)
            files.append(image_path)

    return zip_directory(files, output_zip)


def office_to_pdf(input_file: Path, output_dir: Path) -> Path:
    run_command(
        [
            libreoffice_binary(),
            "--headless",
            "--nologo",
            "--nofirststartwizard",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output_dir),
            str(input_file),
        ],
        timeout=600,
    )
    output_pdf = output_dir / f"{input_file.stem}.pdf"
    if not output_pdf.exists():
        matches = list(output_dir.glob("*.pdf"))
        if not matches:
            raise ConversionError("LibreOffice did not produce a PDF")
        output_pdf = matches[0]
    return output_pdf


def images_to_pdf(image_paths: list[Path], output_pdf: Path) -> Path:
    pil_images: list[Image.Image] = []
    try:
        for path in image_paths:
            image = Image.open(path)
            pil_images.append(image.convert("RGB"))
        if not pil_images:
            raise HTTPException(status_code=400, detail="At least one image is required")
        first, *rest = pil_images
        first.save(output_pdf, save_all=True, append_images=rest)
    finally:
        for image in pil_images:
            image.close()
    return output_pdf


def merge_pdfs(input_pdfs: list[Path], output_pdf: Path) -> Path:
    writer = PdfWriter()
    for pdf_path in input_pdfs:
        reader = PdfReader(str(pdf_path))
        if reader.is_encrypted:
            raise HTTPException(status_code=400, detail=f"{pdf_path.name} is encrypted")
        for page in reader.pages:
            writer.add_page(page)
    with output_pdf.open("wb") as file_obj:
        writer.write(file_obj)
    return output_pdf


def split_pdf(input_pdf: Path, output_zip: Path) -> Path:
    reader = PdfReader(str(input_pdf))
    if reader.is_encrypted:
        raise HTTPException(status_code=400, detail="PDF is encrypted")
    files: list[Path] = []
    for index, page in enumerate(reader.pages, start=1):
        writer = PdfWriter()
        writer.add_page(page)
        page_path = output_zip.with_name(f"page-{index:04d}.pdf")
        with page_path.open("wb") as file_obj:
            writer.write(file_obj)
        files.append(page_path)
    if not files:
        raise HTTPException(status_code=400, detail="PDF has no pages")
    return zip_directory(files, output_zip)


def rotate_pdf(input_pdf: Path, output_pdf: Path, degrees: int) -> Path:
    if degrees not in {90, 180, 270}:
        raise HTTPException(status_code=400, detail="degrees must be one of 90, 180, or 270")

    reader = PdfReader(str(input_pdf))
    if reader.is_encrypted:
        raise HTTPException(status_code=400, detail="PDF is encrypted")
    writer = PdfWriter()
    for page in reader.pages:
        page.rotate(degrees)
        writer.add_page(page)
    with output_pdf.open("wb") as file_obj:
        writer.write(file_obj)
    return output_pdf


def compress_pdf(input_pdf: Path, output_pdf: Path) -> Path:
    run_command(
        [
            ghostscript_binary(),
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            "-dPDFSETTINGS=/ebook",
            "-dNOPAUSE",
            "-dQUIET",
            "-dBATCH",
            f"-sOutputFile={output_pdf}",
            str(input_pdf),
        ],
        timeout=600,
    )
    return output_pdf
